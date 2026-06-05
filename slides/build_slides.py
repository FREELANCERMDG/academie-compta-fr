# -*- coding: utf-8 -*-
"""Genere un diaporama .pptx editable par module (slides/Module-XX.pptx),
derive de la structure markdown des cours. Necessite python-pptx."""
import os, re, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(BASE)
MOD_DIR = os.path.join(ROOT, "modules")

def _hex(h):
    h = (h or "").lstrip("#")
    if len(h) != 6: return None
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

_BR = {"cabinet_name": "", "course_title": "Comptabilité française externalisée",
       "subtitle": "Madagascar - 2026", "primary": "#1F4E78", "accent": "#E8A13A"}
_bp = os.path.join(ROOT, "branding.json")
if os.path.exists(_bp):
    try:
        with open(_bp, encoding="utf-8") as f: _BR.update(json.load(f))
    except Exception: pass

NAVY = _hex(_BR["primary"]) or RGBColor(0x1F, 0x4E, 0x78)
NAVY2 = RGBColor(0x2E, 0x6C, 0xA4)
ACCENT = _hex(_BR["accent"]) or RGBColor(0xE8, 0xA1, 0x3A)
DARK = RGBColor(0x1C, 0x27, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x66, 0x70, 0x7B)
_nm = (_BR.get("cabinet_name") or "").strip()
SUBTITLE = (_nm + "  -  " if _nm else "") + (_BR.get("course_title") or "") + "  -  " + (_BR.get("subtitle") or "")

FILES = sorted([f for f in os.listdir(MOD_DIR) if f.endswith(".md")])

def clean(s):
    s = re.sub(r'`([^`]+)`', r'\1', s)
    s = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', s)
    s = s.replace('**', '').replace('*', '')
    s = re.sub(r'^\[[ xX]\]\s*', '', s)
    # retirer quelques emojis frequents en tete
    s = re.sub(r'^[☀-➿\U0001F000-\U0001FAFF✅✔️⭐⚠⁉‼→←]+\s*', '', s).strip()
    return s.strip()

def parse(md):
    """Retourne (titre, [(section, [bullets])])."""
    lines = md.split('\n')
    title = None
    sections = []
    cur = None
    i = 0
    while i < len(lines):
        line = lines[i]
        st = line.strip()
        h1 = re.match(r'^#\s+(.*)$', st)
        h2 = re.match(r'^##\s+(.*)$', st)
        h3 = re.match(r'^###\s+(.*)$', st)
        if h1 and title is None:
            title = clean(h1.group(1)); i += 1; continue
        if h2:
            cur = {"title": clean(h2.group(1)), "bullets": []}
            sections.append(cur); i += 1; continue
        if cur is None:
            i += 1; continue
        if h3:
            cur["bullets"].append(("sub", clean(h3.group(1)))); i += 1; continue
        # tableau -> resume par une puce
        if '|' in line and i+1 < len(lines) and re.match(r'^\s*\|?\s*:?-{2,}', lines[i+1]):
            hdr = [clean(c) for c in st.strip('|').split('|') if clean(c)]
            cur["bullets"].append(("tab", "Tableau : " + " / ".join(hdr[:5])))
            i += 2
            while i < len(lines) and '|' in lines[i] and lines[i].strip() != '':
                i += 1
            continue
        # liste
        m = re.match(r'^\s*[-*]\s+(.*)$', line)
        if m:
            cur["bullets"].append(("li", clean(m.group(1)))); i += 1; continue
        m = re.match(r'^\s*\d+\.\s+(.*)$', line)
        if m:
            cur["bullets"].append(("li", clean(m.group(1)))); i += 1; continue
        # citation
        if st.startswith('>'):
            cur["bullets"].append(("quote", clean(re.sub(r'^\s*>\s?', '', line)))); i += 1; continue
        if st in ('', '---') or st.startswith('```'):
            i += 1; continue
        # paragraphe court
        p = clean(st)
        if p and len(p) < 160:
            cur["bullets"].append(("li", p))
        i += 1
    return title, sections

# ---- PPTX helpers ----
def band(slide, prs, h=Inches(1.15), color=NAVY):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def title_slide(prs, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.35), prs.slide_width, Inches(0.12))
    acc.fill.solid(); acc.fill.fore_color.rgb = ACCENT; acc.line.fill.background(); acc.shadow.inherit = False
    txt(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(2), title, 40, WHITE, True)
    txt(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.8), SUBTITLE, 16, ACCENT)
    return s

def content_slide(prs, section_title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    band(s, prs)
    txt(s, Inches(0.6), Inches(0.28), Inches(12.1), Inches(0.7), section_title, 24, WHITE, True)
    body = s.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(11.9), Inches(5.6))
    tf = body.text_frame; tf.word_wrap = True
    first = True
    for kind, text in bullets:
        if not text: continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        if kind == "sub":
            r = p.add_run(); r.text = text; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY2
        elif kind == "quote":
            r = p.add_run(); r.text = "  " + text; r.font.size = Pt(14); r.font.italic = True; r.font.color.rgb = GREY
        else:
            r = p.add_run(); r.text = "•  " + (text[:150] + ("…" if len(text) > 150 else "")); r.font.size = Pt(15); r.font.color.rgb = DARK
    return s

def closing_slide(prs, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
    txt(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2), "A retenir", 34, ACCENT, True)
    txt(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(2.5),
        "Appliquez la checklist du module avant toute livraison.\n"
        "Faites le quiz interactif dans la plateforme (seuil 70 %).\n"
        "Vous preparez et proposez ; l'expert-comptable decide et signe.",
        18, WHITE)
    return s

SKIP_KW = ("Evaluation de fin de module", "Évaluation de fin de module")
count = 0
for fn in FILES:
    with open(os.path.join(MOD_DIR, fn), 'r', encoding='utf-8') as f:
        title, sections = parse(f.read())
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    title_slide(prs, title or fn)
    for sec in sections:
        if not sec["bullets"]:
            continue
        if any(k in sec["title"] for k in SKIP_KW):
            content_slide(prs, sec["title"], [("li", "Quiz interactif auto-corrige disponible dans la plateforme e-learning (seuil de reussite : 70 %).")])
            continue
        bl = sec["bullets"]
        # decoupe en pages de 6 puces max
        chunk = []
        page = 1
        for b in bl:
            chunk.append(b)
            if len(chunk) == 6:
                ttl = sec["title"] + (f" (suite {page})" if page > 1 else "")
                content_slide(prs, ttl, chunk); chunk = []; page += 1
        if chunk:
            ttl = sec["title"] + (f" (suite {page})" if page > 1 else "")
            content_slide(prs, ttl, chunk)
    closing_slide(prs, title)
    out = os.path.join(BASE, fn.replace('.md', '.pptx'))
    prs.save(out)
    count += 1
print("Diaporamas generes:", count)
